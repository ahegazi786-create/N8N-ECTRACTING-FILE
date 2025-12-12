from fastapi import FastAPI, UploadFile, File, Header, HTTPException
from pptx import Presentation
import io
import os

app = FastAPI()

API_KEY = os.getenv("API_KEY", "")

def check_key(x_api_key: str | None):
    if API_KEY and x_api_key != API_KEY:
        raise HTTPException(status_code=401, detail="Invalid API key")

@app.post("/pptx-to-json")
async def pptx_to_json(
    file: UploadFile = File(...),
    x_api_key: str | None = Header(default=None),
):
    check_key(x_api_key)

    data = await file.read()
    prs = Presentation(io.BytesIO(data))

    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = " ".join(line.strip() for line in shape.text.splitlines() if line.strip())
                if t:
                    texts.append(t)

        slides.append({
            "slide_index": i,
            "text": "\n".join(texts)
        })

    return {
        "fileName": file.filename,
        "slide_count": len(slides),
        "slides": slides
    }
